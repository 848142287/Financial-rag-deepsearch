"""
PowerPoint文档解析服务
支持提取PPT/PPTX文件的文本、表格和图像信息
"""

import os
import re
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import base64
import io
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    print("python-pptx not available. PPT parsing will be limited.")
    PPTX_AVAILABLE = False

# 尝试导入图像处理库
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    print("PIL not available. Image extraction from PPT will be limited.")
    PIL_AVAILABLE = False

logger = logging.getLogger(__name__)

class PPTParser:
    """PowerPoint文档解析器"""

    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.supported_formats = ['.ppt', '.pptx']

    def can_parse(self, file_path: str) -> bool:
        """检查文件是否为PPT格式"""
        if not PPTX_AVAILABLE:
            return False

        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.supported_formats

    async def parse_ppt(self, file_path: str, extract_images: bool = True) -> Dict[str, Any]:
        """
        解析PPT文件

        Args:
            file_path: PPT文件路径
            extract_images: 是否提取图像

        Returns:
            包含解析结果的字典
        """
        if not self.can_parse(file_path):
            raise ValueError(f"Unsupported file format: {file_path}")

        try:
            self.logger.info(f"开始解析PPT文件: {file_path}")

            # 加载PPT文件
            presentation = Presentation(file_path)

            # 解析结果
            result = {
                'file_path': file_path,
                'file_name': os.path.basename(file_path),
                'total_slides': len(presentation.slides),
                'slides': [],
                'metadata': self._extract_metadata(presentation),
                'text_content': '',
                'tables': [],
                'images': [],
                'charts': [],
                'summary': '',
                'parsing_time': datetime.now().isoformat(),
                'success': True,
                'error': None
            }

            # 解析每一页
            all_text = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_data = self._parse_slide(slide, slide_num, extract_images)
                result['slides'].append(slide_data)

                # 收集所有文本
                if slide_data['text']:
                    all_text.append(slide_data['text'])

                # 收集表格
                if slide_data.get('tables'):
                    result['tables'].extend(slide_data['tables'])

                # 收集图像
                if slide_data.get('images'):
                    result['images'].extend(slide_data['images'])

                # 收集图表
                if slide_data.get('charts'):
                    result['charts'].extend(slide_data['charts'])

            # 生成完整文本内容
            result['text_content'] = '\n\n'.join(all_text)

            # 生成摘要
            result['summary'] = self._generate_summary(result)

            self.logger.info(f"PPT解析完成: {result['total_slides']}页, "
                           f"{len(result['tables'])}个表格, {len(result['images'])}个图像")

            return result

        except Exception as e:
            self.logger.error(f"PPT解析失败: {str(e)}")
            return {
                'file_path': file_path,
                'success': False,
                'error': str(e),
                'parsing_time': datetime.now().isoformat()
            }

    def _extract_metadata(self, presentation) -> Dict[str, Any]:
        """提取PPT元数据"""
        try:
            # 获取文档属性
            core_props = presentation.core_properties

            metadata = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords or '',
                'comments': core_props.comments or '',
                'category': core_props.category or '',
                'created': core_props.created.isoformat() if core_props.created else None,
                'modified': core_props.modified.isoformat() if core_props.modified else None,
                'last_modified_by': core_props.last_modified_by or '',
                'revision': core_props.revision or 0,
                'language': core_props.language or ''
            }

            return metadata

        except Exception as e:
            self.logger.warning(f"提取PPT元数据失败: {str(e)}")
            return {}

    def _parse_slide(self, slide, slide_num: int, extract_images: bool = True) -> Dict[str, Any]:
        """解析单个幻灯片"""
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'text': '',
            'shapes': [],
            'tables': [],
            'images': [],
            'charts': [],
            'notes': '',
            'layout': ''
        }

        try:
            # 获取幻灯片标题
            slide_data['title'] = self._extract_slide_title(slide)

            # 解析所有形状
            text_parts = []
            for shape in slide.shapes:
                shape_data = self._parse_shape(shape, slide_num, extract_images)
                slide_data['shapes'].append(shape_data)

                # 收集文本
                if shape_data.get('text'):
                    text_parts.append(shape_data['text'])

                # 收集表格
                if shape_data.get('table_data'):
                    slide_data['tables'].append(shape_data['table_data'])

                # 收集图像
                if shape_data.get('image_data'):
                    slide_data['images'].append(shape_data['image_data'])

                # 收集图表
                if shape_data.get('chart_type'):
                    slide_data['charts'].append({
                        'slide_number': slide_num,
                        'chart_type': shape_data['chart_type'],
                        'title': shape_data.get('text', ''),
                        'position': shape_data.get('position', {})
                    })

            # 组合幻灯片文本
            slide_data['text'] = '\n\n'.join(text_parts)

            # 获取备注
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text = self._extract_text_from_shape(slide.notes_slide.notes_text_frame)
                slide_data['notes'] = notes_text

            # 获取布局信息
            slide_data['layout'] = getattr(slide.slide_layout, 'name', 'Unknown')

        except Exception as e:
            self.logger.error(f"解析幻灯片{slide_num}失败: {str(e)}")
            slide_data['error'] = str(e)

        return slide_data

    def _extract_slide_title(self, slide) -> str:
        """提取幻灯片标题"""
        try:
            # 优先查找标题占位符
            for shape in slide.shapes:
                if shape.is_placeholder and hasattr(shape.placeholder_format, 'type'):
                    if shape.placeholder_format.type == 1:  # 标题占位符
                        return self._extract_text_from_shape(shape)

            # 如果没有找到标题占位符，查找第一个文本框
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = self._extract_text_from_shape(shape)
                    if text and len(text.strip()) > 0:
                        # 通常标题比较短，取前100个字符
                        return text[:100].strip()

            return ''

        except Exception as e:
            self.logger.warning(f"提取幻灯片标题失败: {str(e)}")
            return ''

    def _parse_shape(self, shape, slide_num: int, extract_images: bool = True) -> Dict[str, Any]:
        """解析形状"""
        shape_data = {
            'shape_id': shape.shape_id,
            'shape_type': str(shape.shape_type),
            'name': shape.name,
            'text': '',
            'position': {},
            'size': {},
            'format': {},
            'table_data': None,
            'image_data': None,
            'chart_type': None
        }

        try:
            # 获取位置和大小
            if hasattr(shape, 'left'):
                shape_data['position'] = {
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }

            # 根据形状类型解析
            if hasattr(shape, 'text_frame') and shape.text_frame:
                shape_data['text'] = self._extract_text_from_shape(shape)
                shape_data['format'] = self._extract_text_format(shape)

            elif hasattr(shape, 'table'):
                shape_data['table_data'] = self._extract_table_data(shape.table)
                shape_data['shape_type'] = 'TABLE'

            elif hasattr(shape, 'chart'):
                shape_data['chart_type'] = self._identify_chart_type(shape.chart)
                shape_data['shape_type'] = 'CHART'

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and extract_images:
                shape_data['image_data'] = self._extract_image_data(shape, slide_num)
                shape_data['shape_type'] = 'PICTURE'

            # 处理组合形状
            elif hasattr(shape, 'shapes'):
                shape_data['shapes'] = []
                for sub_shape in shape.shapes:
                    sub_data = self._parse_shape(sub_shape, slide_num, extract_images)
                    if sub_data.get('text') or sub_data.get('image_data'):
                        shape_data['shapes'].append(sub_data)

        except Exception as e:
            self.logger.warning(f"解析形状失败: {str(e)}")
            shape_data['error'] = str(e)

        return shape_data

    def _extract_text_from_shape(self, shape) -> str:
        """从形状中提取文本"""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return ''

            text_frame = shape.text_frame
            paragraphs = []

            for paragraph in text_frame.paragraphs:
                paragraph_text = ""
                for run in paragraph.runs:
                    if run.text:
                        paragraph_text += run.text

                if paragraph_text.strip():
                    paragraphs.append(paragraph_text.strip())

            return '\n'.join(paragraphs)

        except Exception as e:
            self.logger.warning(f"提取形状文本失败: {str(e)}")
            return ''

    def _extract_text_format(self, shape) -> Dict[str, Any]:
        """提取文本格式信息"""
        format_info = {}

        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame = shape.text_frame

                # 提取字体信息
                fonts = []
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font:
                            font_info = {
                                'name': run.font.name,
                                'size': run.font.size.pt if run.font.size else None,
                                'bold': run.font.bold,
                                'italic': run.font.italic,
                                'color': self._color_to_hex(run.font.color) if run.font.color else None
                            }
                            fonts.append(font_info)

                if fonts:
                    format_info['fonts'] = fonts

                # 提取段落格式
                paragraphs_info = []
                for paragraph in text_frame.paragraphs:
                    para_info = {
                        'alignment': str(paragraph.alignment) if paragraph.alignment else None,
                        'level': paragraph.level
                    }
                    paragraphs_info.append(para_info)

                if paragraphs_info:
                    format_info['paragraphs'] = paragraphs_info

        except Exception as e:
            self.logger.warning(f"提取文本格式失败: {str(e)}")

        return format_info

    def _extract_table_data(self, table) -> Dict[str, Any]:
        """提取表格数据"""
        try:
            table_data = {
                'rows': table.rows.__len__(),
                'columns': table.columns.__len__(),
                'data': []
            }

            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = self._extract_text_from_shape(cell)
                    row_data.append(cell_text.strip())
                table_data['data'].append(row_data)

            return table_data

        except Exception as e:
            self.logger.warning(f"提取表格数据失败: {str(e)}")
            return None

    def _extract_image_data(self, shape, slide_num: int) -> Dict[str, Any]:
        """提取图像数据"""
        try:
            if not PIL_AVAILABLE:
                return {'error': 'PIL not available for image processing'}

            # 获取图像
            image_bytes = shape.image.blob

            # 创建PIL图像对象
            image = Image.open(io.BytesIO(image_bytes))

            # 转换为base64编码
            buffered = io.BytesIO()
            image.save(buffered, format="PNG")
            image_base64 = base64.b64encode(buffered.getvalue()).decode()

            image_data = {
                'slide_number': slide_num,
                'width': image.width,
                'height': image.height,
                'format': image.format,
                'mode': image.mode,
                'size_bytes': len(image_bytes),
                'base64_data': image_base64,
                'position': shape_data.get('position', {}) if hasattr(self, 'shape_data') else {}
            }

            return image_data

        except Exception as e:
            self.logger.warning(f"提取图像数据失败: {str(e)}")
            return {'error': str(e)}

    def _identify_chart_type(self, chart) -> str:
        """识别图表类型"""
        try:
            # 简化的图表类型识别
            chart_type = str(type(chart).__name__)

            # 根据图表名称进行更具体的判断
            if 'Bar' in chart_type:
                return '柱状图'
            elif 'Line' in chart_type:
                return '折线图'
            elif 'Pie' in chart_type:
                return '饼图'
            elif 'Area' in chart_type:
                return '面积图'
            elif 'Scatter' in chart_type:
                return '散点图'
            else:
                return '图表'

        except Exception as e:
            self.logger.warning(f"识别图表类型失败: {str(e)}")
            return '未知图表'

    def _color_to_hex(self, color) -> Optional[str]:
        """将颜色转换为十六进制"""
        try:
            if hasattr(color, 'rgb'):
                rgb = color.rgb
                return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
        except:
            pass
        return None

    def _generate_summary(self, result: Dict[str, Any]) -> str:
        """生成PPT内容摘要"""
        try:
            summary_parts = []

            # 基本信息
            file_info = f"PPT文档《{result['file_name']}》共{result['total_slides']}页"
            summary_parts.append(file_info)

            # 主要内容
            if result['text_content']:
                # 取前500个字符作为内容预览
                content_preview = result['text_content'][:500]
                if len(result['text_content']) > 500:
                    content_preview += "..."
                summary_parts.append(f"主要内容：{content_preview}")

            # 结构信息
            structure_info = []
            if result['tables']:
                structure_info.append(f"{len(result['tables'])}个表格")
            if result['images']:
                structure_info.append(f"{len(result['images'])}个图像")
            if result['charts']:
                structure_info.append(f"{len(result['charts'])}个图表")

            if structure_info:
                summary_parts.append(f"包含：{', '.join(structure_info)}")

            # 元数据
            metadata = result['metadata']
            if metadata.get('title'):
                summary_parts.append(f"标题：{metadata['title']}")
            if metadata.get('author'):
                summary_parts.append(f"作者：{metadata['author']}")

            return '\n'.join(summary_parts)

        except Exception as e:
            self.logger.warning(f"生成PPT摘要失败: {str(e)}")
            return f"PPT文档《{result.get('file_name', 'Unknown')}》，共{result.get('total_slides', 0)}页"

# 创建全局PPT解析器实例
ppt_parser = PPTParser()

def parse_ppt_file(file_path: str, extract_images: bool = True) -> Dict[str, Any]:
    """
    解析PPT文件的便捷函数

    Args:
        file_path: PPT文件路径
        extract_images: 是否提取图像

    Returns:
        解析结果字典
    """
    import asyncio
    return asyncio.run(ppt_parser.parse_ppt(file_path, extract_images))