"""
文档转换服务 - 支持 Word/PPT/Excel 转 PDF

功能特点：
1. Word (.docx, .doc) → PDF
2. PowerPoint (.pptx, .ppt) → PDF
3. Excel (.xlsx, .xls) → PDF
4. 纯文本 (.txt) → PDF
5. 统一的转换接口
"""

import os
import io
import tempfile
from typing import Optional, List
from pathlib import Path
from dataclasses import dataclass

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.units import inch
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False


@dataclass
class ConversionResult:
    """转换结果"""
    success: bool
    output_path: Optional[str] = None
    output_bytes: Optional[bytes] = None
    page_count: int = 0
    error_message: str = ""
    metadata: dict = None

    def __post_init__(self):
        if self.metadata is None:
            self.metadata = {}


class DocumentConverter:
    """
    文档转换器

    支持将各种文档格式转换为 PDF
    """

    def __init__(self, output_dir: Optional[str] = None):
        """
        初始化转换器

        Args:
            output_dir: 输出目录，如果为None则使用临时目录
        """
        self.output_dir = output_dir or tempfile.gettempdir()
        self._ensure_output_dir()

    def _ensure_output_dir(self):
        """确保输出目录存在"""
        os.makedirs(self.output_dir, exist_ok=True)

    def convert_to_pdf(
        self,
        file_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        自动检测文件类型并转换为PDF

        Args:
            file_path: 输入文件路径
            output_path: 输出PDF路径（可选）

        Returns:
            ConversionResult对象
        """
        if not os.path.exists(file_path):
            return ConversionResult(
                success=False,
                error_message=f"文件不存在: {file_path}"
            )

        file_ext = Path(file_path).suffix.lower()

        # 根据文件扩展名选择转换方法
        if file_ext == '.pdf':
            return self._convert_pdf_to_pdf(file_path, output_path)
        elif file_ext in ['.docx', '.doc']:
            return self.convert_word_to_pdf(file_path, output_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self.convert_ppt_to_pdf(file_path, output_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self.convert_excel_to_pdf(file_path, output_path)
        elif file_ext in ['.txt', '.md']:
            return self.convert_text_to_pdf(file_path, output_path)
        elif file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif']:
            return self.convert_image_to_pdf(file_path, output_path)
        else:
            return ConversionResult(
                success=False,
                error_message=f"不支持的文件格式: {file_ext}"
            )

    def convert_word_to_pdf(
        self,
        word_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        Word 转 PDF

        Args:
            word_path: Word文件路径
            output_path: 输出PDF路径

        Returns:
            ConversionResult对象
        """
        if not DOCX_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="python-docx未安装，请安装: pip install python-docx"
            )

        if not PYMUPDF_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="PyMuPDF未安装，请安装: pip install pymupdf"
            )

        try:
            # 读取Word文档
            doc = DocxDocument(word_path)

            # 创建临时PDF
            if output_path is None:
                output_path = os.path.join(
                    self.output_dir,
                    f"{Path(word_path).stem}.pdf"
                )

            # 使用PyMuPDF创建PDF
            pdf_doc = fitz.open()

            # 简单文本提取（实际应用中需要更复杂的排版）
            for para in doc.paragraphs:
                if para.text.strip():
                    page = pdf_doc.new_page()
                    page.insert_text(
                        (50, 700),
                        para.text,
                        fontsize=12,
                        fontname="helvetica"
                    )

            pdf_doc.save(output_path)
            pdf_doc.close()

            return ConversionResult(
                success=True,
                output_path=output_path,
                page_count=len(doc.paragraphs)
            )

        except Exception as e:
            return ConversionResult(
                success=False,
                error_message=f"Word转换失败: {str(e)}"
            )

    def convert_ppt_to_pdf(
        self,
        ppt_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        PowerPoint 转 PDF

        Args:
            ppt_path: PowerPoint文件路径
            output_path: 输出PDF路径

        Returns:
            ConversionResult对象
        """
        if not PPTX_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="python-pptx未安装，请安装: pip install python-pptx"
            )

        if not PYMUPDF_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="PyMuPDF未安装，请安装: pip install pymupdf"
            )

        try:
            # 读取PowerPoint文档
            prs = PptxPresentation(ppt_path)

            # 创建临时PDF
            if output_path is None:
                output_path = os.path.join(
                    self.output_dir,
                    f"{Path(ppt_path).stem}.pdf"
                )

            # 使用PyMuPDF创建PDF
            pdf_doc = fitz.open()

            # 提取每页内容
            for slide_num, slide in enumerate(prs.slides):
                page = pdf_doc.new_page()

                # 提取文本和形状
                y_position = 750
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        page.insert_text(
                            (50, y_position),
                            shape.text.strip(),
                            fontsize=12,
                            fontname="helvetica"
                        )
                        y_position -= 20

                        if y_position < 50:
                            # 换页
                            page = pdf_doc.new_page()
                            y_position = 750

            pdf_doc.save(output_path)
            pdf_doc.close()

            return ConversionResult(
                success=True,
                output_path=output_path,
                page_count=len(prs.slides)
            )

        except Exception as e:
            return ConversionResult(
                success=False,
                error_message=f"PPT转换失败: {str(e)}"
            )

    def convert_excel_to_pdf(
        self,
        excel_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        Excel 转 PDF

        Args:
            excel_path: Excel文件路径
            output_path: 输出PDF路径

        Returns:
            ConversionResult对象
        """
        # Excel转换需要更复杂的处理，这里提供基础实现
        return ConversionResult(
            success=False,
            error_message="Excel转PDF功能暂未实现，建议先将Excel另存为PDF"
        )

    def convert_text_to_pdf(
        self,
        text_path: str,
        output_path: Optional[str] = None,
        font_size: int = 12
    ) -> ConversionResult:
        """
        纯文本/TXT 转 PDF

        Args:
            text_path: 文本文件路径
            output_path: 输出PDF路径
            font_size: 字体大小

        Returns:
            ConversionResult对象
        """
        if not REPORTLAB_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="reportlab未安装，请安装: pip install reportlab"
            )

        try:
            # 读取文本
            with open(text_path, 'r', encoding='utf-8') as f:
                text = f.read()

            # 创建输出路径
            if output_path is None:
                output_path = os.path.join(
                    self.output_dir,
                    f"{Path(text_path).stem}.pdf"
                )

            # 创建PDF
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet

            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            # 分割文本为段落
            paragraphs = text.split('\n\n')
            for para_text in paragraphs:
                if para_text.strip():
                    para = Paragraph(para_text.strip(), styles['Normal'])
                    story.append(para)

            doc.build(story)

            return ConversionResult(
                success=True,
                output_path=output_path,
                page_count=1  # 简化处理
            )

        except Exception as e:
            return ConversionResult(
                success=False,
                error_message=f"文本转换失败: {str(e)}"
            )

    def convert_image_to_pdf(
        self,
        image_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        图片转PDF

        Args:
            image_path: 图片文件路径
            output_path: 输出PDF路径

        Returns:
            ConversionResult对象
        """
        if not PYMUPDF_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="PyMuPDF未安装，请安装: pip install pymupdf"
            )

        try:
            from PIL import Image

            # 打开图片
            img = Image.open(image_path)

            # 创建输出路径
            if output_path is None:
                output_path = os.path.join(
                    self.output_dir,
                    f"{Path(image_path).stem}.pdf"
                )

            # 转换为RGB
            if img.mode != "RGB":
                img = img.convert("RGB")

            # 保存为PDF
            img.save(output_path, "PDF", resolution=100.0)

            return ConversionResult(
                success=True,
                output_path=output_path,
                page_count=1
            )

        except Exception as e:
            return ConversionResult(
                success=False,
                error_message=f"图片转换失败: {str(e)}"
            )

    def _convert_pdf_to_pdf(
        self,
        pdf_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        PDF 到 PDF（直接复制）

        Args:
            pdf_path: PDF文件路径
            output_path: 输出PDF路径

        Returns:
            ConversionResult对象
        """
        if not PYMUPDF_AVAILABLE:
            return ConversionResult(
                success=False,
                error_message="PyMuPDF未安装"
            )

        try:
            if output_path is None:
                output_path = os.path.join(
                    self.output_dir,
                    f"{Path(pdf_path).stem}.pdf"
                )

            # 直接复制
            import shutil
            shutil.copy2(pdf_path, output_path)

            # 获取页数
            doc = fitz.open(output_path)
            page_count = doc.page_count
            doc.close()

            return ConversionResult(
                success=True,
                output_path=output_path,
                page_count=page_count
            )

        except Exception as e:
            return ConversionResult(
                success=False,
                error_message=f"PDF复制失败: {str(e)}"
            )


# 创建全局实例
_document_converter = None


def get_document_converter() -> DocumentConverter:
    """获取文档转换器实例"""
    global _document_converter
    if _document_converter is None:
        _document_converter = DocumentConverter()
    return _document_converter
