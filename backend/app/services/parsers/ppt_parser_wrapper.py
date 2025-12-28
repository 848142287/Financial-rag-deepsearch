"""
增强PowerPoint文档解析器 - 支持多模态分析、结构化输出
继承BaseFileParser，提供标准化的PPT/PPTX解析功能
"""

import logging
import os
import tempfile
import shutil
from typing import Dict, List, Any, Optional
from pathlib import Path
import asyncio
from dataclasses import dataclass, field
from datetime import datetime

from .base import BaseFileParser, ParseResult
from .multimodal_analyzer import get_multimodal_analyzer, ImageAnalysisResult

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available")


@dataclass
class SlideElement:
    """幻灯片元素"""
    element_type: str  # 'text', 'image', 'chart', 'table', 'formula'
    content: str  # 内容或描述
    metadata: Dict[str, Any] = field(default_factory=dict)
    position: int = 0


@dataclass
class SlideSection:
    """幻灯片（作为一节）"""
    slide_number: int
    title: str
    elements: List[SlideElement] = field(default_factory=list)
    layout: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)


class PPTParserWrapper(BaseFileParser):
    """增强PowerPoint文档解析器 - 支持多模态分析"""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
        self.enable_multimodal = self.config.get('enable_multimodal', True)
        self.extract_images = self.config.get('extract_images', True)
        self.enable_ocr = self.config.get('enable_ocr', True)

        # 初始化多模态分析器
        if self.enable_multimodal:
            self.multimodal_analyzer = get_multimodal_analyzer(self.config)
        else:
            self.multimodal_analyzer = None

    @property
    def supported_extensions(self) -> List[str]:
        """支持的文件扩展名列表"""
        return ['.pptx', '.ppt']

    @property
    def parser_name(self) -> str:
        """解析器名称"""
        return "EnhancedPPTParser"

    def can_parse(self, file_path: str, file_extension: str = None) -> bool:
        """检查是否能解析指定文件"""
        if not PPTX_AVAILABLE:
            return False

        if file_extension:
            return file_extension.lower() in self.supported_extensions

        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.supported_extensions

    async def parse(self, file_path: str, **kwargs) -> ParseResult:
        """
        解析PPT文件（增强版 - 支持多模态分析）

        Args:
            file_path: PPT文件路径
            **kwargs: 其他参数
                - extract_images: 是否提取图像 (默认True)
                - enable_ocr: 是否启用OCR (默认True)
                - output_format: 输出格式 ('markdown', 'json', 默认'markdown')

        Returns:
            ParseResult: 解析结果
        """
        import time
        start_time = time.time()
        temp_dir = None

        try:
            if not PPTX_AVAILABLE:
                return ParseResult(
                    content="",
                    metadata={'file_type': 'pptx', 'error': 'python-pptx library not available'},
                    success=False,
                    error_message='python-pptx library not installed'
                )

            # 获取参数
            extract_images = kwargs.get('extract_images', self.extract_images)
            enable_ocr = kwargs.get('enable_ocr', self.enable_ocr)
            output_format = kwargs.get('output_format', 'markdown')

            # 创建临时目录用于提取图片
            if extract_images:
                temp_dir = tempfile.mkdtemp(prefix='pptx_images_')

            # 1. 解析PPT结构
            prs = Presentation(file_path)
            slides_data = await self._parse_presentation(prs, temp_dir)

            # 2. 多模态分析
            if extract_images and self.multimodal_analyzer:
                await self._analyze_slides_multimodal(slides_data, temp_dir, enable_ocr)

            # 3. 生成输出
            if output_format == 'markdown':
                content = self._generate_markdown(slides_data)
            else:
                content = self._generate_json_output(slides_data)

            # 4. 构建元数据
            metadata = self._build_metadata_enhanced(file_path, prs, slides_data)
            metadata['parse_time'] = time.time() - start_time

            # 5. 清理临时目录
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

            return ParseResult(
                content=content,
                metadata=metadata,
                success=True,
                parse_time=time.time() - start_time
            )

        except Exception as e:
            logger.error(f"PPT解析失败: {e}", exc_info=True)

            # 清理临时目录
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

            return ParseResult(
                content="",
                metadata={'file_type': 'pptx', 'error': str(e)},
                success=False,
                error_message=str(e),
                parse_time=time.time() - start_time
            )

    async def _parse_presentation(
        self,
        prs: Presentation,
        temp_dir: Optional[str]
    ) -> List[SlideSection]:
        """解析演示文稿结构"""
        slides_data = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_section = SlideSection(
                slide_number=slide_idx,
                title=self._extract_slide_title(slide),
                layout=self._get_slide_layout_name(slide)
            )

            element_position = 0

            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 处理文本框
                if hasattr(shape, "text") and shape.text.strip():
                    slide_section.elements.append(SlideElement(
                        element_type='text',
                        content=shape.text.strip(),
                        position=element_position
                    ))
                    element_position += 1

                # 处理图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_info = await self._extract_image(shape, temp_dir, slide_idx, element_position)
                    if image_info:
                        slide_section.elements.append(SlideElement(
                            element_type='image',
                            content="",
                            metadata=image_info,
                            position=element_position
                        ))
                        element_position += 1

                # 处理表格
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_info = self._extract_table(shape)
                    slide_section.elements.append(SlideElement(
                        element_type='table',
                        content="",
                        metadata=table_info,
                        position=element_position
                    ))
                    element_position += 1

                # 处理图表（需要特殊处理，因为pptx库对图表的支持有限）
                if hasattr(shape, "chart"):
                    chart_info = self._extract_chart(shape)
                    if chart_info:
                        slide_section.elements.append(SlideElement(
                            element_type='chart',
                            content="",
                            metadata=chart_info,
                            position=element_position
                        ))
                        element_position += 1

            slides_data.append(slide_section)

        return slides_data

    async def _extract_image(
        self,
        shape,
        temp_dir: Optional[str],
        slide_number: int,
        position: int
    ) -> Optional[Dict[str, Any]]:
        """提取图片"""
        if not temp_dir:
            return None

        try:
            # 获取图片数据
            image = shape.image
            ext = image.ext
            image_bytes = image.blob

            # 保存到临时文件
            image_filename = f"slide_{slide_number}_img_{position}.{ext}"
            image_path = os.path.join(temp_dir, image_filename)

            with open(image_path, 'wb') as f:
                f.write(image_bytes)

            return {
                'image_type': 'unknown',  # 将由多模态分析器确定
                'temp_path': image_path,
                'filename': image_filename,
                'content_type': image.content_type
            }

        except Exception as e:
            logger.error(f"Failed to extract image: {e}")
            return None

    def _extract_table(self, shape) -> Dict[str, Any]:
        """提取表格"""
        try:
            table = shape.table
            table_data = []

            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)

            return {
                'rows': len(table_data),
                'columns': len(table_data[0]) if table_data else 0,
                'data': table_data
            }

        except Exception as e:
            logger.error(f"Failed to extract table: {e}")
            return {}

    def _extract_chart(self, shape) -> Optional[Dict[str, Any]]:
        """提取图表"""
        try:
            chart = shape.chart
            return {
                'chart_type': str(chart.chart_type),
                'has_title': chart.has_title,
                'title': chart.chart_title.text_frame.text if chart.has_title else ""
            }
        except:
            return None

    def _extract_slide_title(self, slide) -> str:
        """提取幻灯片标题"""
        try:
            # 查找标题占位符
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                    return shape.text.strip()
                elif shape.has_text_frame:
                    text = shape.text.strip()
                    # 简单启发式：第一个简短的文本可能是标题
                    if len(text) < 100 and text:
                        return text
            return f"幻灯片 {slide.slides.index(slide) + 1}"
        except:
            return ""

    def _get_slide_layout_name(self, slide) -> str:
        """获取幻灯片布局名称"""
        try:
            return slide.slide_layout.name
        except:
            return "未知布局"

    async def _analyze_slides_multimodal(
        self,
        slides_data: List[SlideSection],
        temp_dir: Optional[str],
        enable_ocr: bool
    ):
        """多模态分析幻灯片"""
        if not self.multimodal_analyzer:
            return

        for slide in slides_data:
            for element in slide.elements:
                if element.element_type == 'image':
                    image_path = element.metadata.get('temp_path')
                    if image_path and os.path.exists(image_path):
                        # 使用多模态分析器分析图片
                        result = await self.multimodal_analyzer.analyze_image(
                            image_path,
                            image_type='auto'
                        )

                        # 更新元素
                        element.content = result.description
                        element.metadata.update({
                            'ocr_text': result.ocr_text,
                            'image_type': result.image_type,
                            'chart_info': result.chart_info,
                            'formula_info': result.formula_info,
                            'confidence': result.confidence
                        })

                        # 根据图片类型更新元素类型
                        if result.image_type == 'chart':
                            element.element_type = 'chart'
                        elif result.image_type == 'formula':
                            element.element_type = 'formula'

    def _generate_markdown(self, slides_data: List[SlideSection]) -> str:
        """生成Markdown输出"""
        markdown_parts = []

        for slide in slides_data:
            # 幻灯片标题
            markdown_parts.append(f"\n## {slide.title} (第{slide.slide_number}页)\n")

            # 幻灯片元素
            for element in slide.elements:
                if element.element_type == 'text':
                    markdown_parts.append(f"{element.content}\n")
                elif element.element_type == 'image':
                    img_desc = element.content or "图片"
                    ocr_text = element.metadata.get('ocr_text', '')
                    markdown_parts.append(f"\n**图片**: {img_desc}\n")
                    if ocr_text:
                        markdown_parts.append(f"*图片中的文字*: {ocr_text}\n")
                elif element.element_type == 'chart':
                    chart_desc = element.content or "图表"
                    chart_info = element.metadata.get('chart_info', {})
                    markdown_parts.append(f"\n**图表**:\n\n{chart_desc}\n\n")
                    if chart_info.get('x_axis'):
                        markdown_parts.append(f"- *横坐标*: {chart_info['x_axis']}\n")
                    if chart_info.get('y_axis'):
                        markdown_parts.append(f"- *纵坐标*: {chart_info['y_axis']}\n")
                    if chart_info.get('trend'):
                        markdown_parts.append(f"- *趋势*: {chart_info['trend']}\n")
                elif element.element_type == 'table':
                    table_data = element.metadata.get('data', [])
                    markdown_parts.append(f"\n{self._format_table_as_markdown(table_data)}\n")
                elif element.element_type == 'formula':
                    formula_desc = element.content or "公式"
                    markdown_parts.append(f"\n**公式**: {formula_desc}\n")

        return "\n".join(markdown_parts)

    def _format_table_as_markdown(self, table_data: List[List[str]]) -> str:
        """将表格格式化为Markdown"""
        if not table_data:
            return ""

        # 确保第一行不是None
        first_row = table_data[0] or []
        if not first_row:
            return ""

        # 表头
        header = "| " + " | ".join(str(cell) if cell is not None else "" for cell in first_row) + " |"
        separator = "| " + " | ".join(["---"] * len(first_row)) + " |"

        # 表体
        rows = []
        for row in table_data[1:]:
            # 确保row不是None
            safe_row = row if row is not None else []
            rows.append("| " + " | ".join(str(cell) if cell is not None else "" for cell in safe_row) + " |")

        return "\n".join([header, separator] + rows)

    def _generate_json_output(self, slides_data: List[SlideSection]) -> str:
        """生成JSON输出"""
        import json

        output = {
            'slides': [
                {
                    'slide_number': slide.slide_number,
                    'title': slide.title,
                    'layout': slide.layout,
                    'elements': [
                        {
                            'type': el.element_type,
                            'content': el.content,
                            'metadata': el.metadata,
                            'position': el.position
                        }
                        for el in slide.elements
                    ],
                    'metadata': slide.metadata
                }
                for slide in slides_data
            ]
        }

        return json.dumps(output, ensure_ascii=False, indent=2)

    def _build_metadata_enhanced(
        self,
        file_path: str,
        prs: Presentation,
        slides_data: List[SlideSection]
    ) -> Dict[str, Any]:
        """构建增强元数据"""
        from pptx.opc.constants import CONTENT_TYPE as CT

        # 统计元素数量
        element_counts = {}
        for slide in slides_data:
            for element in slide.elements:
                element_counts[element.element_type] = element_counts.get(element.element_type, 0) + 1

        metadata = {
            'file_type': 'pptx',
            'parser_version': '2.0',
            'multimodal_enabled': self.enable_multimodal,
            'ocr_enabled': self.enable_ocr,
            'total_slides': len(slides_data),
            'element_counts': element_counts,
            'texts': element_counts.get('text', 0),
            'images': element_counts.get('image', 0),
            'charts': element_counts.get('chart', 0),
            'tables': element_counts.get('table', 0),
            'formulas': element_counts.get('formula', 0),
        }

        # 添加文档属性
        try:
            core_props = prs.core_properties
            metadata.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'comments': core_props.comments or ''
            })
        except:
            pass

        return metadata
